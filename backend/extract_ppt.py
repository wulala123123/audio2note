# filename: backend/extract_ppt.py
import cv2
import tqdm
from pathlib import Path
from skimage.metrics import structural_similarity as ssim
from pptx import Presentation
from pptx.util import Inches

def extract_key_frames_persistent_reference(video_path: Path, output_base_dir: Path, ssim_threshold=0.9, frame_interval=15, stability_frames=5, stability_threshold=0.995):
    """
    从视频中提取关键帧，采用“持久参考点”逻辑。
    只有在新的关键帧被确认后，才会更新用于变化检测的参考帧。
    """
    video_path = Path(video_path)
    output_base_dir = Path(output_base_dir)

    # 1. 检查和设置路径
    if not video_path.exists():
        print(f"错误：视频文件不存在于 '{video_path}'")
        return None

    # video_name 将被用于文件命名 - 修改为从父目录获取名称
    # Example: output/video/my_video/my_video_cropped.mp4 -> my_video
    video_name = video_path.parent.name
    image_output_dir = output_base_dir / 'ppt_images' / video_name
    ppt_output_dir = output_base_dir / 'pptx_files'
    ppt_output_path = ppt_output_dir / f"{video_name}.pptx"

    # 新增：检查PPT是否已存在
    if ppt_output_path.exists():
        print(f"\n跳过: PPT文件 '{ppt_output_path}' 已存在。")
        return ppt_output_path

    image_output_dir.mkdir(parents=True, exist_ok=True)
    ppt_output_dir.mkdir(parents=True, exist_ok=True)
    
    print(f"\n{'='*20} 开始提取: {video_name} {'='*20}")
    print(f"图片将保存到: {image_output_dir}")
    print(f"PPT将保存到: {ppt_output_path}")

    # 2. 初始化
    cap = cv2.VideoCapture(str(video_path))
    if not cap.isOpened():
        print(f"错误：无法打开视频文件 '{video_path}'")
        return None

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    last_saved_gray = None
    saved_frame_count = 0
    current_frame_index = -1
    candidate_frame = None
    stable_counter = 0
    candidate_frame_index = -1
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))

    # 3. 循环处理视频帧
    with tqdm.tqdm(total=total_frames, desc=f"提取 {video_name}") as pbar:
        while cap.isOpened():
            ret, frame = cap.read()
            if not ret:
                break
            current_frame_index += 1
            pbar.update(1)
            
            if current_frame_index > 0 and current_frame_index % frame_interval != 0:
                continue

            current_frame_gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)

            # 逻辑1：处理第一帧
            if last_saved_gray is None:
                img_path = image_output_dir / f"{video_name}_{saved_frame_count:04d}.jpg"
                cv2.imwrite(str(img_path), frame)
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(str(img_path), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
                saved_frame_count += 1
                last_saved_gray = current_frame_gray
                continue

            # --- 第二级过滤器：变化与稳定性检测 ---
            score, _ = ssim(last_saved_gray, current_frame_gray, full=True)

            if score < ssim_threshold:
                if candidate_frame is None:
                    candidate_frame = frame.copy()
                    candidate_frame_index = current_frame_index
                    stable_counter = 1
                else:
                    candidate_frame_gray = cv2.cvtColor(candidate_frame, cv2.COLOR_BGR2GRAY)
                    stability_score, _ = ssim(candidate_frame_gray, current_frame_gray, full=True)
                    
                    if stability_score > stability_threshold:
                        stable_counter += 1
                    else:
                        candidate_frame = frame.copy()
                        candidate_frame_index = current_frame_index
                        stable_counter = 1
            else:
                if candidate_frame is not None:
                    candidate_frame = None
                    stable_counter = 0

            # 逻辑3：确认并保存稳定的候选帧
            if candidate_frame is not None and stable_counter >= stability_frames:
                img_path = image_output_dir / f"{video_name}_{saved_frame_count:04d}.jpg"
                cv2.imwrite(str(img_path), candidate_frame)
                
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(str(img_path), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
                
                saved_frame_count += 1
                last_saved_gray = cv2.cvtColor(candidate_frame, cv2.COLOR_BGR2GRAY)
                candidate_frame = None
                stable_counter = 0

    # 4. 保存PPT并释放资源
    if saved_frame_count > 0:
        prs.save(str(ppt_output_path))
        print(f"\n处理完成！共提取 {saved_frame_count} 帧关键帧。")
        print(f"演示文稿已保存至: {ppt_output_path}")
        result_path = ppt_output_path
    else:
        print("\n未提取到任何关键帧。请尝试调整参数。")
        result_path = None
        
    cap.release()
    cv2.destroyAllWindows()
    return result_path

def main():
    """主执行函数，用于查找并处理所有已裁剪的视频。"""
    # 自动定位项目根目录
    current_dir = Path(__file__).resolve().parent
    project_root = current_dir.parent
    
    # 假设输出目录还是在 root/output
    output_base_dir = project_root / 'output'
    
    # --- 1. 查找所有由 crop_ppt.py 生成的已裁剪视频 ---
    cropped_video_base_dir = output_base_dir / 'video'
    # search_pattern: output/video/*/*_cropped.mp4
    cropped_video_paths = list(cropped_video_base_dir.glob('*/*_cropped.mp4'))

    if not cropped_video_paths:
        print(f"错误：在 '{cropped_video_base_dir}' 目录中没有找到已裁剪的视频。")
        print(f"请先运行 crop_ppt.py 来生成裁剪后的视频。")
        return
    
    print(f"找到 {len(cropped_video_paths)} 个已裁剪的视频，准备开始提取PPT...")

    # --- 2. 循环处理每个视频 ---
    for video_file_path in cropped_video_paths:
        ## --- 参数调优指南 --- ##
        extract_key_frames_persistent_reference(
            video_path=video_file_path,
            output_base_dir=output_base_dir,
            ssim_threshold=0.95,
            frame_interval=20,
            stability_frames=3,
            stability_threshold=0.995
        )
    
    print(f"\n{'='*25} 所有视频提取完毕 {'='*25}")


if __name__ == '__main__':
    main()