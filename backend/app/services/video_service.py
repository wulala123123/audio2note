"""
文件名: video_service.py
功能描述: 视频处理核心服务，实现 GPU 加速的 PPT 提取与音频转录编排
核心逻辑:
    - _locate_ppt_region(): 使用 Canny 边缘检测定位视频中的 PPT 区域
    - _generate_lightweight_video(): 生成轻量视频 (640px, 5fps) 用于加速分析
    - _run_funnel_analysis(): 三层漏斗 PPT 提取 (L1帧差 + L2清晰度 + L3 OCR去重)
    - _high_res_capture(): 高清回溯 - 从原视频截取最终画面
    - process(): 主入口，编排 PPT 提取与音频转录两个独立模块

全链路架构 (Lightweight Media Workflow):
    1. Step 1.1: ROI Detection - 定位 PPT 区域
    2. Step 1.2: Lightweight Video - 生成轻量视频 (640px, 5fps, 去音频)
    3. Step 1.3: Funnel Analysis - 三层漏斗分析轻量视频
       - L1: 帧差检测 (场景分割)
       - L2: 清晰度择优 (选冠军帧)
       - L3: OCR 语义去重 (过滤重复页 + 非PPT页面)
    4. Step 1.4: High-Res Capture - 从原视频高清回溯

设计亮点:
    - **Timestamp First**: 所有逻辑基于时间戳 (秒 float)，严禁依赖 frame_index
    - 轻量视频分析 (快速) + 原视频截取 (高清) 分离
    - Generator 模式流式输出，支持实时进度更新
    - 流程结束自动清理临时文件 (轻量视频)
"""
import cv2
import shutil
from pathlib import Path
from typing import Tuple, Optional

from pptx import Presentation
from pptx.util import Inches
from loguru import logger

from app.core.config import OUTPUT_DIR, TEMP_DIR
from app.core.task_manager import update_task_progress
from app.services.audio_service import get_audio_transcriber
from app.services.gpu_frame_processor import GPUFrameProcessor, BestShot
from app.services.ocr_deduper import OCRDeduper
from app.utils.ffmpeg_utils import (
    generate_lightweight_video,
    extract_frame_at_timestamp,
    extract_frames_batch
)


class VideoService:
    """
    视频处理服务主类
    
    职责: 编排整个视频 -> PPT 转换流程，协调各子模块工作
    
    核心流程 (Lightweight Media Workflow):
        1. 定位 PPT 区域 (ROI Detection)
        2. 生成轻量视频 (640px, 5fps)
        3. 在轻量视频上运行三层漏斗分析
        4. 用时间戳回溯原视频截取高清画面
    
    Attributes:
        output_guid: 任务唯一标识，用于组织输出目录
        base_output_path: 任务输出根目录
        frame_processor: GPU 帧处理器 (L1 + L2)
        ocr_deduper: OCR 语义去重器 (L3)
    
    Example:
        >>> service = VideoService(output_guid="task-123")
        >>> result = service.process(Path("lecture.mp4"), enable_ppt_extraction=True)
        >>> print(result["ppt_file"])
    """
    
    def __init__(self, output_guid: str) -> None:
        """
        初始化视频处理服务
        
        Args:
            output_guid: 任务唯一标识符 (通常为 UUID)
        """
        self.output_guid = output_guid
        self.base_output_path = OUTPUT_DIR / output_guid
        
        # 定义子目录结构
        # 轻量视频临时目录：放入 temp 下，流程结束后自动清理
        self.temp_video_dir = TEMP_DIR / output_guid
        self.debug_images_dir = self.base_output_path / "debug_images"
        self.ppt_images_dir = self.base_output_path / "ppt_images"
        self.ppt_output_dir = self.base_output_path / "ppt_output"
        self.transcripts_dir = self.base_output_path / "transcripts"
        
        # 创建所需文件夹
        for p in [self.temp_video_dir, self.debug_images_dir, 
                  self.ppt_images_dir, self.ppt_output_dir, self.transcripts_dir]:
            p.mkdir(parents=True, exist_ok=True)
        
        logger.debug(f"📁 输出目录已创建: {self.base_output_path}")
        logger.debug(f"📁 临时目录已创建: {self.temp_video_dir}")
        
        # ========== 初始化 GPU 处理器 (L1 + L2) ==========
        # 参数说明:
        #   diff_threshold: 帧间差异阈值，超过此值视为场景切换
        #   min_scene_duration: 场景最短持续时间，过滤动态视频片段
        #   sample_interval: 采样间隔 (秒)，每 0.2 秒取一次样 (每秒 5 个点)
        self.frame_processor = GPUFrameProcessor(
            diff_threshold=0.05,
            min_scene_duration=1,
            sample_interval=0.2  # 每 0.2 秒采样一次
        )
        
        # ========== 初始化 OCR 去重器 (L3) ==========
        # 参数说明:
        #   similarity_threshold: 文本相似度阈值，超过则判定为重复页
        self.ocr_deduper = OCRDeduper(
            similarity_threshold=0.90
        )

    def process(
        self, 
        input_video_path: Path, 
        enable_ppt_extraction: bool = True,
        enable_audio_transcription: bool = True
    ) -> dict:
        """
        视频处理主入口: 编排 PPT 提取与音频转录两个独立模块
        
        两个功能模块完全解耦，可独立启用或同时启用。
        进度条会根据启用的模块数量自动分配区间。
        
        Args:
            input_video_path: 原始视频文件路径
            enable_ppt_extraction: 是否执行 PPT 提取流程 (默认 True)
            enable_audio_transcription: 是否执行音频转录流程 (默认 True)
            
        Returns:
            dict: 处理结果，包含各输出文件路径
                - guid: 任务 ID
                - ppt_file: PPT 文件路径 (如启用 PPT 提取)
                - transcript_file: 转录文件路径 (如启用音频转录)
        
        Raises:
            ValueError: PPT 区域定位失败或视频裁剪失败
        """
        input_video_path = Path(input_video_path)
        
        logger.info("=" * 50)
        logger.info(f"🎬 VideoService.process() 开始处理")
        logger.info(f"   📂 输入: {input_video_path.name}")
        logger.info(f"   🆔 GUID: {self.output_guid}")
        logger.info(f"   📊 PPT提取: {enable_ppt_extraction} | 🎤 音频转录: {enable_audio_transcription}")
        logger.info("=" * 50)
        
        ppt_path = None
        transcript_path = None
        
        # ============================================================
        #               模块 1: PPT 提取 (条件执行)
        # ============================================================
        if enable_ppt_extraction:
            logger.info("📊 [PPT 提取模块] 开始执行 (Lightweight Media Workflow)...")
            
            # 进度区间分配:
            #   - 若同时启用音频: PPT 占 0-85%, 音频占 85-100%
            #   - 若仅 PPT: PPT 占 0-100%
            ppt_progress_end = 85 if enable_audio_transcription else 100
            
            # ----- Step 1.1: 定位 PPT 区域 -----
            update_task_progress(self.output_guid, 5, "正在定位 PPT 区域...")
            logger.info("🔍 Step 1.1: 定位 PPT 区域 (Canny 边缘检测)")
            
            bbox = self._locate_ppt_region(input_video_path)
            
            if not bbox:
                logger.error("❌ 无法定位 PPT 区域")
                raise ValueError("无法定位 PPT 区域，请确保视频中包含清晰的 PPT 画面")
            
            logger.success(f"✅ PPT 区域定位成功: x={bbox[0]}, y={bbox[1]}, w={bbox[2]}, h={bbox[3]}")
            
            # ----- Step 1.2: 生成轻量视频 -----
            update_task_progress(self.output_guid, 10, "正在生成轻量视频 (GPU 加速)...")
            logger.info("🎥 Step 1.2: 生成轻量视频 (640px, 5fps)")
            
            lightweight_video_path = self._generate_lightweight_video(input_video_path, bbox)
            
            if not lightweight_video_path:
                logger.error("❌ 轻量视频生成失败")
                raise ValueError("轻量视频生成失败")
            
            logger.success(f"✅ 轻量视频生成完成: {lightweight_video_path.name}")
            
            # ----- Step 1.3: 三层漏斗分析 -----
            update_task_progress(self.output_guid, 25, "正在进行三层漏斗分析...")
            logger.info("🎯 Step 1.3: 三层漏斗分析 (L1→L2→L3)")
            
            final_timestamps = self._run_funnel_analysis(lightweight_video_path)
            
            logger.info(f"📊 漏斗分析结果: 共 {len(final_timestamps)} 个有效时间点")
            
            if not final_timestamps:
                logger.warning("⚠️ 未检测到任何有效 PPT 页面")
                ppt_path = None
            else:
                # ----- Step 1.4: 高清回溯 -----
                update_task_progress(self.output_guid, 70, "正在高清回溯截取...")
                logger.info("📸 Step 1.4: 高清回溯 (从原视频截取)")
                
                ppt_path = self._high_res_capture(
                    source_video=input_video_path,
                    timestamps=final_timestamps,
                    crop_bbox=bbox
                )
                
                if ppt_path:
                    logger.success(f"✅ PPT 生成完成: {ppt_path.name}")
                else:
                    logger.warning("⚠️ PPT 生成失败")
        
        # ============================================================
        #               模块 2: 音频转录 (条件执行，完全独立)
        # ============================================================
        if enable_audio_transcription:
            logger.info("🎤 [音频转录模块] 开始执行...")
            
            # 进度区间:
            #   - 若同时启用 PPT: 从 85% 开始
            #   - 若仅音频: 从 0% 开始
            audio_progress_start = 85 if enable_ppt_extraction else 0
            
            update_task_progress(
                self.output_guid, 
                audio_progress_start + 5, 
                "正在进行语音识别 (FunASR)..."
            )
            
            try:
                logger.info("🔊 调用 FunASR 进行本地语音识别...")
                transcript_text = get_audio_transcriber().transcribe_video(input_video_path)
                
                if transcript_text:
                    transcript_path = self.transcripts_dir / f"{self.output_guid}.txt"
                    with open(transcript_path, "w", encoding="utf-8") as f:
                        f.write(transcript_text)
                    logger.success(f"✅ 转录文件已保存: {transcript_path.name}")
                    logger.debug(f"   📝 转录内容预览: {transcript_text[:100]}...")
                else:
                    logger.warning("⚠️ 转录结果为空")
                    
            except Exception as e:
                logger.exception(f"❌ 音频转录过程出错: {e}")
        
        # ============================================================
        #               流程结束: 清理临时文件
        # ============================================================
        self._cleanup_temp_files()
        
        # ========== 返回结果 ==========
        result = {
            "guid": self.output_guid,
            "ppt_file": str(ppt_path) if ppt_path else None,
            "transcript_file": str(transcript_path) if transcript_path else None
        }
        
        update_task_progress(self.output_guid, 100, "处理完成")
        
        logger.info("=" * 50)
        logger.info(f"🏁 VideoService.process() 处理完成")
        logger.info(f"   📄 PPT: {'✅ ' + ppt_path.name if ppt_path else '❌ 未生成'}")
        logger.info(f"   📝 转录: {'✅ ' + transcript_path.name if transcript_path else '❌ 未生成'}")
        logger.info("=" * 50)
        
        return result

    def _cleanup_temp_files(self) -> None:
        """
        清理临时文件
        
        在处理流程结束后调用，删除轻量视频等临时文件。
        """
        try:
            if self.temp_video_dir.exists():
                shutil.rmtree(self.temp_video_dir)
                logger.info(f"🗑️ 已清理临时目录: {self.temp_video_dir}")
        except Exception as e:
            logger.warning(f"⚠️ 清理临时目录失败: {e}")

    def _locate_ppt_region(self, video_path: Path) -> Tuple[int, int, int, int] | None:
        """
        定位视频中的 PPT 区域 (边缘检测法)
        
        算法策略:
            1. 在视频 20%/40%/60% 位置各采样一帧
            2. 使用 Canny 边缘检测识别边缘
            3. 使用轮廓分析寻找最大四边形区域
            4. 返回该区域的 bounding box
        
        Why 多点采样?
            - 视频开头/结尾可能没有 PPT 画面
            - 多点采样提高检测成功率
        
        Args:
            video_path: 输入视频路径
            
        Returns:
            tuple: (x, y, w, h) PPT 区域坐标和尺寸，失败返回 None
        """
        logger.debug(f"🔍 开始定位 PPT 区域: {video_path.name}")
        
        cap = cv2.VideoCapture(str(video_path))
        if not cap.isOpened():
            logger.error(f"❌ 无法打开视频: {video_path}")
            return None
        
        try:
            total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            sample_points = [0.2, 0.4, 0.6]  # 采样点: 20%, 40%, 60%
            
            logger.debug(f"   📊 总帧数: {total_frames}, 采样点: {sample_points}")
            
            for point in sample_points:
                frame_idx = int(total_frames * point)
                cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
                ret, frame = cap.read()
                
                if not ret:
                    logger.warning(f"   ⚠️ 采样点 {point:.0%} 读取失败")
                    continue
                
                logger.debug(f"   🖼️ 分析采样点 {point:.0%} (帧 {frame_idx})")
                
                # 保存调试图像 (可视化边缘检测过程)
                cv2.imwrite(str(self.debug_images_dir / "0_original.jpg"), frame)
                
                # ----- Canny 边缘检测流水线 -----
                # Step 1: BGR -> Gray (减少计算量)
                gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
                cv2.imwrite(str(self.debug_images_dir / "1_gray.jpg"), gray)
                
                # Step 2: 高斯模糊 (去噪，平滑边缘)
                blurred = cv2.GaussianBlur(gray, (5, 5), 0)
                
                # Step 3: Canny 边缘检测
                # Why (30, 120)? 低阈值 30 检测弱边缘，高阈值 120 过滤噪点
                edged = cv2.Canny(blurred, 30, 120)
                cv2.imwrite(str(self.debug_images_dir / "2_edged.jpg"), edged)
                
                # ----- 轮廓分析 -----
                contours, _ = cv2.findContours(
                    edged.copy(), 
                    cv2.RETR_EXTERNAL,      # 只检测外轮廓
                    cv2.CHAIN_APPROX_SIMPLE  # 压缩轮廓点
                )
                
                if not contours:
                    logger.debug(f"   ⚠️ 采样点 {point:.0%} 未找到轮廓")
                    continue
                
                # 取面积最大的 5 个轮廓 (PPT 通常是最大的矩形区域)
                contours = sorted(contours, key=cv2.contourArea, reverse=True)[:5]
                
                for c in contours:
                    # 轮廓近似: 减少顶点数量
                    peri = cv2.arcLength(c, True)
                    approx = cv2.approxPolyDP(c, 0.03 * peri, True)
                    
                    # 筛选条件:
                    #   1. 必须是 4 边形 (PPT 是矩形)
                    #   2. 面积占比 > 10% (过滤小区域)
                    frame_area = frame.shape[0] * frame.shape[1]
                    area_ratio = cv2.contourArea(c) / frame_area
                    
                    if len(approx) == 4 and area_ratio > 0.1:
                        # 保存调试结果
                        debug_img = frame.copy()
                        cv2.drawContours(debug_img, [approx], -1, (0, 255, 0), 3)
                        cv2.imwrite(str(self.debug_images_dir / "3_final_region.jpg"), debug_img)
                        
                        bbox = cv2.boundingRect(approx)
                        logger.info(f"   ✅ 在采样点 {point:.0%} 找到 PPT 区域")
                        logger.info(f"      📐 Bounding Box: x={bbox[0]}, y={bbox[1]}, w={bbox[2]}, h={bbox[3]}")
                        logger.info(f"      📊 面积占比: {area_ratio:.1%}")
                        return bbox
            
            logger.error("❌ 所有采样点均未找到有效 PPT 区域")
            return None
            
        finally:
            cap.release()

    def _generate_lightweight_video(
        self, 
        source_video: Path, 
        crop_bbox: Tuple[int, int, int, int]
    ) -> Optional[Path]:
        """
        生成轻量视频 (核心优化)
        
        调用 FFmpeg 生成低分辨率轻量视频用于后续分析。
        
        轻量视频参数:
            - crop: 只保留 PPT 区域
            - scale: 宽缩放到 640px
            - fps: 降帧到 5 FPS
            - audio: 去除音频
        
        Args:
            source_video: 原始视频路径
            crop_bbox: 裁剪区域 (x, y, w, h)
        
        Returns:
            Path: 轻量视频路径，失败返回 None
        """
        lightweight_path = self.temp_video_dir / f"{self.output_guid}_lightweight.mp4"
        
        def progress_callback(percent: int, message: str) -> None:
            """轻量视频生成进度回调 (占 10-25%)"""
            actual_progress = 10 + int(percent * 0.15)
            update_task_progress(self.output_guid, actual_progress, message)
        
        result = generate_lightweight_video(
            source_video=source_video,
            output_path=lightweight_path,
            crop_box=crop_bbox,
            target_width=640,
            target_fps=5,
            progress_callback=progress_callback
        )
        
        return result

    def _run_funnel_analysis(self, lightweight_video: Path) -> list[float]:
        """
        三层漏斗分析 (运行在轻量视频上)
        
        在轻量视频上执行 L1+L2+L3 分析，输出最终时间戳列表。
        
        处理流程:
            L1 (物理层): GPU 帧差检测 → 场景分割
            L2 (质量层): 拉普拉斯清晰度 → 选冠军帧
            L3 (语义层): OCR 文本识别 → 过滤重复页 + 非PPT页面
        
        关键设计:
            - 所有逻辑基于时间戳 (秒 float)
            - 使用轻量视频进行 OCR (快速)
            - 无文字内容的帧视为非 PPT 页面，自动过滤
        
        Args:
            lightweight_video: 轻量视频路径 (640px, 5fps)
        
        Returns:
            list[float]: 最终时间戳列表，如 [1.2, 15.6, 48.2, ...]
        """
        logger.info("🔄 开始三层漏斗分析...")
        logger.info("   L1: GPU 帧差检测 → 场景分割")
        logger.info("   L2: 拉普拉斯清晰度 → 选冠军帧")
        logger.info("   L3: OCR 识别 → 过滤重复页 + 非PPT页面")
        
        # 重置 OCR 去重器
        self.ocr_deduper.reset()
        
        final_timestamps: list[float] = []
        candidate_count = 0
        
        # ----- L1 + L2: GPU 帧差 + 清晰度择优 -----
        def l1l2_progress(percent: int, message: str) -> None:
            """L1+L2 进度回调 (占 25-50%)"""
            actual_progress = 25 + int(percent * 0.25)
            update_task_progress(self.output_guid, actual_progress, message)
        
        for best_shot in self.frame_processor.extract_best_shots(
            lightweight_video, 
            progress_callback=l1l2_progress
        ):
            candidate_count += 1
            
            logger.debug(f"   🎬 候选帧 #{candidate_count}: "
                        f"timestamp={best_shot.timestamp:.2f}s, "
                        f"清晰度={best_shot.sharpness_score:.4f}")
            
            # ----- L3: OCR 识别与过滤 -----
            update_task_progress(
                self.output_guid, 
                50 + int((candidate_count / max(candidate_count, 1)) * 20),
                f"L3 OCR 分析: 第 {candidate_count} 个候选"
            )
            
            # 从轻量视频读取帧进行 OCR (轻量视频足够进行文字识别)
            frame = self.frame_processor.get_frame_at_timestamp(
                lightweight_video, 
                best_shot.timestamp
            )
            
            if frame is None:
                logger.warning(f"   ⚠️ 无法读取帧 @ {best_shot.timestamp:.2f}s")
                continue
            
            is_duplicate, text = self.ocr_deduper.is_duplicate(frame)
            
            # 过滤条件 1: 无文字内容 → 非 PPT 页面
            if not text or not text.strip():
                logger.debug(f"   📄 @ {best_shot.timestamp:.2f}s 无文字内容，判定为非PPT页面，跳过")
                continue
            
            # 过滤条件 2: 与已保存页面重复
            if is_duplicate:
                logger.debug(f"   🔄 @ {best_shot.timestamp:.2f}s 与已保存页相似度过高，跳过")
                continue
            
            # 保留该时间戳
            final_timestamps.append(best_shot.timestamp)
            self.ocr_deduper.mark_as_saved(text)
            
            logger.info(f"   ✅ 保留: @ {best_shot.timestamp:.2f}s (第 {len(final_timestamps)} 页)")
        
        logger.success(f"✅ 漏斗分析完成: {candidate_count} 候选 → {len(final_timestamps)} 保留")
        return final_timestamps

    def _high_res_capture(
        self, 
        source_video: Path,
        timestamps: list[float],
        crop_bbox: Tuple[int, int, int, int]
    ) -> Optional[Path]:
        """
        高清回溯: 从原视频截取最终画面并生成 PPTX
        
        遍历时间戳列表，使用 FFmpeg 从原视频精确截取高清帧，
        然后组装成 PPTX 文件。
        
        关键设计:
            - 从原视频截取 (保留完整分辨率)
            - 使用原始 bbox 坐标裁剪
            - 不缩放，保持最高画质
        
        Args:
            source_video: 原始视频路径
            timestamps: 最终时间戳列表 (秒)
            crop_bbox: 裁剪区域 (x, y, w, h)
        
        Returns:
            Path: 生成的 PPTX 文件路径，失败返回 None
        """
        if not timestamps:
            logger.warning("⚠️ 时间戳列表为空，无法生成 PPT")
            return None
        
        logger.info(f"📸 开始高清回溯: 共 {len(timestamps)} 个时间点")
        
        # ----- 批量截取高清帧 -----
        def capture_progress(percent: int, message: str) -> None:
            """截取进度回调 (占 70-90%)"""
            actual_progress = 70 + int(percent * 0.2)
            update_task_progress(self.output_guid, actual_progress, message)
        
        frame_paths = extract_frames_batch(
            source_video=source_video,
            timestamps=timestamps,
            output_dir=self.ppt_images_dir,
            crop_box=None,  # 不裁剪，保留完整原视频画面
            progress_callback=capture_progress
        )
        
        if not frame_paths:
            logger.warning("⚠️ 未能截取任何帧")
            return None
        
        # ----- 组装 PPTX -----
        update_task_progress(self.output_guid, 92, "正在生成 PPTX...")
        logger.info(f"📄 组装 PPTX: {len(frame_paths)} 页")
        
        ppt_path = self.ppt_output_dir / f"{self.output_guid}.pptx"
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        
        for i, img_path in enumerate(frame_paths):
            # 添加空白幻灯片
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # 添加图片 (全屏)
            slide.shapes.add_picture(
                str(img_path),
                Inches(0), 
                Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )
            
            logger.debug(f"   📄 添加第 {i+1} 页: {img_path.name}")
        
        prs.save(str(ppt_path))
        logger.success(f"✅ PPTX 生成完成: {ppt_path.name} ({len(frame_paths)} 页)")
        
        return ppt_path
