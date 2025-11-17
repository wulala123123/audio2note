import cv2
import os
from skimage.metrics import structural_similarity as ssim
from pptx import Presentation
from pptx.util import Inches
import glob
import tqdm

def extract_key_frames_persistent_reference(video_path, output_base_dir='output', ssim_threshold=0.9, frame_interval=15, stability_frames=5, stability_threshold=0.995):
    """
    从视频中提取关键帧，采用“持久参考点”逻辑。
    只有在新的关键帧被确认后，才会更新用于变化检测的参考帧。
    """
    # 1. 检查和设置路径
    if not os.path.exists(video_path):
        print(f"错误：视频文件不存在于 '{video_path}'")
        return
    # video_name 将被用于文件命名 - 修改为从父目录获取名称
    video_name = os.path.basename(os.path.dirname(video_path))
    image_output_dir = os.path.join(output_base_dir, 'ppt_images', video_name)
    ppt_output_dir = os.path.join(output_base_dir, 'pptx_files')
    ppt_output_path = os.path.join(ppt_output_dir, f"{video_name}.pptx")

    # 新增：检查PPT是否已存在
    if os.path.exists(ppt_output_path):
        print(f"\n跳过: PPT文件 '{ppt_output_path}' 已存在。")
        return

    os.makedirs(image_output_dir, exist_ok=True)
    os.makedirs(ppt_output_dir, exist_ok=True)
    
    print(f"\n{'='*20} 开始提取: {video_name} {'='*20}")
    print(f"图片将保存到: {image_output_dir}")
    print(f"PPT将保存到: {ppt_output_path}")

    # 2. 初始化
    cap = cv2.VideoCapture(video_path)
    if not cap.isOpened():
        print(f"错误：无法打开视频文件 '{video_path}'")
        return
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    last_saved_gray = None
    saved_frame_count = 0
    current_frame_index = -1
    candidate_frame = None
    stable_counter = 0
    candidate_frame_index = -1
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))

    # 3. 循环处理视频帧
    with tqdm.tqdm(total=total_frames, desc=f"提取 {video_name}") as pbar:
        while cap.isOpened():
            ret, frame = cap.read()
            if not ret:
                break
            current_frame_index += 1
            pbar.update(1)
            
            if current_frame_index > 0 and current_frame_index % frame_interval != 0:
                continue

            current_frame_gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)

            # 逻辑1：处理第一帧
            if last_saved_gray is None:
                img_path = os.path.join(image_output_dir, f"{video_name}_{saved_frame_count:04d}.jpg")
                cv2.imwrite(img_path, frame)
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
                saved_frame_count += 1
                last_saved_gray = current_frame_gray
                continue

            # --- 第二级过滤器：变化与稳定性检测 ---
            score, _ = ssim(last_saved_gray, current_frame_gray, full=True)

            if score < ssim_threshold:
                if candidate_frame is None:
                    candidate_frame = frame.copy()
                    candidate_frame_index = current_frame_index
                    stable_counter = 1
                else:
                    candidate_frame_gray = cv2.cvtColor(candidate_frame, cv2.COLOR_BGR2GRAY)
                    stability_score, _ = ssim(candidate_frame_gray, current_frame_gray, full=True)
                    
                    if stability_score > stability_threshold:
                        stable_counter += 1
                    else:
                        candidate_frame = frame.copy()
                        candidate_frame_index = current_frame_index
                        stable_counter = 1
            else:
                if candidate_frame is not None:
                    candidate_frame = None
                    stable_counter = 0

            # 逻辑3：确认并保存稳定的候选帧
            if candidate_frame is not None and stable_counter >= stability_frames:
                img_path = os.path.join(image_output_dir, f"{video_name}_{saved_frame_count:04d}.jpg")
                cv2.imwrite(img_path, candidate_frame)
                
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
                
                saved_frame_count += 1
                last_saved_gray = cv2.cvtColor(candidate_frame, cv2.COLOR_BGR2GRAY)
                candidate_frame = None
                stable_counter = 0

    # 4. 保存PPT并释放资源
    if saved_frame_count > 0:
        prs.save(ppt_output_path)
        print(f"\n处理完成！共提取 {saved_frame_count} 帧关键帧。")
        print(f"演示文稿已保存至: {ppt_output_path}")
    else:
        print("\n未提取到任何关键帧。请尝试调整参数。")
    cap.release()
    cv2.destroyAllWindows()

def main():
    """主执行函数，用于查找并处理所有已裁剪的视频。"""
    # --- 1. 查找所有由 crop_ppt.py 生成的已裁剪视频 ---
    cropped_video_base_dir = os.path.join('output', 'video')
    search_pattern = os.path.join(cropped_video_base_dir, '*', '*_cropped.mp4')
    cropped_video_paths = glob.glob(search_pattern)

    if not cropped_video_paths:
        print(f"错误：在 '{cropped_video_base_dir}' 目录中没有找到已裁剪的视频。")
        print(f"请先运行 crop_ppt.py 来生成裁剪后的视频。")
        return
    
    print(f"找到 {len(cropped_video_paths)} 个已裁剪的视频，准备开始提取PPT...")

    # --- 2. 循环处理每个视频 ---
    for video_file_path in cropped_video_paths:
        ## --- 参数调优指南 --- ##
        # 下面是四个核心参数，你可以通过调整它们来适应不同类型的视频。
        extract_key_frames_persistent_reference(
            video_path=video_file_path,
            output_base_dir='output',

            # 1. SSIM 变化检测阈值 (ssim_threshold)
            ssim_threshold=0.95,

            # 2. 帧检查间隔 (frame_interval)
            frame_interval=20,

            # 3. 画面稳定确认数 (stability_frames)
            stability_frames=3,

            # 4. 稳定性判断阈值 (stability_threshold)
            stability_threshold=0.995
        )
    
    print(f"\n{'='*25} 所有视频提取完毕 {'='*25}")


if __name__ == '__main__':
    main()